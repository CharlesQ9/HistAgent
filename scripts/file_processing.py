from smolagents import Tool
# import easyocr
import os
import re
import json
import time
import torch
import requests
import pdfminer.high_level
import mammoth
import pandas as pd
import pptx
from smolagents.models import MessageRole, Model
import base64
from pathlib import Path
from typing import Optional, Dict, Any, List, Union
import subprocess
from translate import Translator
import puremagic
import tempfile
import shutil
import tabulate

class FileProcessor:
    """Base class for file processing, providing shared resources and common functionality"""
    
    def __init__(
        self,
        ocr_languages=["en", "ch_sim"],  # Use EasyOCR supported language codes
        model=None  # Add model parameter
    ):
        self.ocr_reader = None
        self.ocr_languages = ocr_languages
        self.model = model

    def detect_file_type(self, file_path: str) -> str:
        """Detect file type and return extension"""
        # Try to get extension from file name first
        ext = os.path.splitext(file_path)[1].lower()
        
        # If no extension, use puremagic to detect
        if not ext:
            try:
                guesses = puremagic.magic_file(file_path)
                if guesses:
                    ext = "." + guesses[0].extension.strip()
            except (FileNotFoundError, IsADirectoryError, PermissionError):
                pass
                
        return ext
    def get_appropriate_tool(self, file_path: str) -> Optional[Tool]:
        """Return the appropriate processing tool based on file type"""
        ext = self.detect_file_type(file_path)
        
        # Mapping of file types to tools
        tool_map = {
            '.pdf': PDFTool(self),
            '.docx': DOCXTool(self),
            '.xlsx': XLSXTool(self),
            '.xls': XLSXTool(self),
            '.pptx': PPTXTool(self)
        }
        
        return tool_map.get(ext)

    def get_ocr_reader(self):
        """Initialize OCR reader with the correct language codes"""
        if self.ocr_reader is None:
            try:
                import easyocr
                self.ocr_reader = easyocr.Reader(self.ocr_languages, gpu=False)
            except Exception as e:
                raise Exception(f"OCR reader initialization failed: {str(e)}")
        return self.ocr_reader

class OCRTool(Tool):
    name = "OCR_Tool"
    description = "Extract text from images using OCR with automatic language detection."
    inputs = {
        "image_path": {
            "type": "string",
            "description": "Path to the image file.",
            "nullable": True
        },
        "languages": {
            "type": "array",
            "items": {
                "type": "string"
            },
            "description": "Optional list of language codes for OCR. If not provided, will auto-detect.",
            "default": ["en", "ch_sim"],
            "nullable": True
        },
        "auto_detect": {
            "type": "boolean",
            "description": "Whether to automatically detect languages in the image.",
            "default": True,
            "nullable": True
        }
    }
    output_type = "string"

    def __init__(self, file_processor: FileProcessor, model=None):
        super().__init__()
        self.file_processor = file_processor
        self.model = model  # LLM model for result optimization
        self._easyocr_readers = {}  # Cache readers for different language combinations
        
        # Common language combinations and their corresponding codes
        self.language_groups = {
            "cjk": ["ch_sim", "ch_tra", "ja", "ko", "en"],  # Chinese + Japanese + Korean + English
            "european": ["en", "fr", "de", "es", "it", "pt", "ru"],  # European languages
            "indic": ["hi", "ta", "te", "kn", "mr", "ne", "en"],  # Indic languages + English
            "arabic": ["ar", "fa", "ur", "en"],  # Arabic + English
            "default": ["en", "ch_sim"]  # Default combination
        }

    def _get_reader(self, languages=None):
        """Get or create EasyOCR reader"""
        if languages is None:
            languages = self.language_groups["default"]
            
        # Convert language list to sorted tuple for dictionary key
        lang_key = tuple(sorted(languages))
        
        # If there's a cached reader, reuse it
        if lang_key in self._easyocr_readers:
            return self._easyocr_readers[lang_key]
            
        # Otherwise create a new reader
        try:
            import easyocr
            print(f"Initializing EasyOCR, supported languages: {languages}")
            # Modify: Explicitly specify GPU=False to avoid CUDA-related issues
            reader = easyocr.Reader(languages, gpu=False)
            self._easyocr_readers[lang_key] = reader
            return reader
        except Exception as e:
            print(f"Failed to initialize EasyOCR: {str(e)}")
            # Modify: Return default reader instead of None
            if lang_key != tuple(sorted(self.language_groups["default"])):
                print("Attempting to use default language combination...")
                return self._get_reader(self.language_groups["default"])
            return None

    def _detect_language(self, image_path):
        """Use image analysis to detect the main language in the image"""
        try:
            # First try using basic English OCR to get some text samples
            basic_reader = self._get_reader(["en"])
            if not basic_reader:
                return self.language_groups["default"]
                
            sample_text = basic_reader.readtext(image_path, detail=0)
            
            # If no text is detected, return default language combination
            if not sample_text:
                return self.language_groups["default"]
                
            # Analyze text features to guess language
            text = " ".join(sample_text)
            
            # Detect CJK characters (Chinese, Japanese, Korean)
            if any(ord(c) > 0x4E00 and ord(c) < 0x9FFF for c in text):
                return self.language_groups["cjk"]
                
            # Detect Cyrillic characters (Russian, etc.)
            if any(ord(c) > 0x0400 and ord(c) < 0x04FF for c in text):
                return ["ru", "en"]
                
            # Detect Arabic characters
            if any(ord(c) > 0x0600 and ord(c) < 0x06FF for c in text):
                return self.language_groups["arabic"]
                
            # Detect Indic languages
            if any(ord(c) > 0x0900 and ord(c) < 0x097F for c in text):
                return self.language_groups["indic"]
                
            # If no special characters, assume it's European language
            return self.language_groups["european"]
            
        except Exception as e:
            print(f"Language detection failed: {str(e)}")
            return self.language_groups["default"]

    def forward(self, image_path: str = None, languages: List[str] = None, auto_detect: bool = True) -> str:
        try:
            if image_path is None:
                return "Error: No image file path provided"
                
            if not os.path.exists(image_path):
                return f"Error: File path not found: {image_path}"
                
            # Check file extension
            ext = os.path.splitext(image_path)[1].lower()
            if ext not in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.webp']:
                return "Error: Unsupported image format. Supported formats include: JPG, JPEG, PNG, BMP, TIFF, WEBP"
            
            # If no languages provided and auto-detection is enabled, detect language
            if languages is None and auto_detect:
                print("Detecting image language...")
                languages = self._detect_language(image_path)
                print(f"Detected possible language combination: {languages}")
            elif languages is None:
                languages = self.language_groups["default"]
                
            # Get OCR reader
            reader = self._get_reader(languages)
            if reader is None:
                # Modify: Return more detailed error info and try system description
                print("Failed to initialize OCR engine, trying system tool to describe image...")
                if self.model:
                    return self._describe_image_with_llm(image_path)
                return "Error: Failed to initialize OCR engine, please check EasyOCR installation and language support"
                
            # Execute OCR
            print(f"Executing OCR on image {image_path}, using languages: {languages}")
            try:
                # Add timeout handling
                results = reader.readtext(image_path, detail=0)  # detail=0只返回文本
            except Exception as ocr_error:
                print(f"OCR processing failed: {str(ocr_error)}")
                # Try using default language
                if languages != self.language_groups["default"]:
                    print("Trying default language combination...")
                    reader = self._get_reader(self.language_groups["default"])
                    if reader:
                        try:
                            results = reader.readtext(image_path, detail=0)
                        except:
                            results = []
                    else:
                        results = []
                else:
                    results = []
            
            # Check results
            if not results:
                print("OCR did not detect text, try using other language combinations....")
                # Try using other language combinations
                for group_name, group_langs in self.language_groups.items():
                    if group_langs != languages:
                        print(f"Trying {group_name} language combination: {group_langs}")
                        alt_reader = self._get_reader(group_langs)
                        if alt_reader:
                            try:
                                alt_results = alt_reader.readtext(image_path, detail=0)
                                if alt_results:
                                    results = alt_results
                                    print(f"Successfully detected text using {group_name} language combination")
                                    break
                            except Exception as e:
                                print(f"Failed to use {group_name} language combination: {str(e)}")
                
                # If still no results, try using LLM to describe image
                if not results and self.model:
                    return self._describe_image_with_llm(image_path)
                elif not results:
                    return "No text detected from image, please try using other tools or manual processing"
                
            # Combine results into text
            text = "\n".join(results)
            
            # Optimize OCR results with LLM (if model is available)
            if self.model and self._needs_optimization(text):
                return self._optimize_with_llm(text, image_path)
            
            return text
            
        except Exception as e:
            print(f"OCR processing failed: {str(e)}")
            # If model is available, try using LLM to describe image
            if self.model:
                return self._describe_image_with_llm(image_path)
            return f"OCR processing failed: {str(e)}"

    def _needs_optimization(self, text: str) -> bool:
        """Determine whether the OCR results need optimization."""
        if not text:
            return False
            
        # Check if text contains too many special characters
        special_chars = sum(1 for c in text if not c.isalnum() and not c.isspace())
        if special_chars / len(text) > 0.3:  # If special characters exceed 30%
            return True
            
        # Check if there are too many consecutive non-word characters (It might be a recognition error.)
        if re.search(r'[^\w\s]{3,}', text):
            return True
            
        # Check if there are too many uncommon character combinations
        if re.search(r'[a-z][A-Z]{2,}|[A-Z][a-z]{2,}[A-Z]', text):
            return True
            
        return True
        
    def _optimize_with_llm(self, text: str, image_path: str) -> str:
        """Use LLM to optimize OCR results"""
        try:
            # Use the same message format as TextInspectorTool
            messages = [
                {
                    "role": MessageRole.SYSTEM,
                    "content": [
                        {
                            "type": "text",
                            "text": "你是一个OCR文本修正专家。你的任务是修复OCR识别错误，提供准确的文本。只返回修正后的文本，不要添加任何解释、分析或评论。"
                        }
                    ],
                },
                {
                    "role": MessageRole.USER,
                    "content": [
                        {
                            "type": "text",
                            "text": f"以下是OCR识别的原始文本，可能包含错误:\n\n{text}\n\n请直接提供修正后的纯文本，无需任何额外内容:"
                        }
                    ],
                }
            ]
            
            print("DEBUG: Using the same call method as TextInspectorTool")
            
            # 直接调用模型并获取content属性
            response = self.model(messages).content
            
            # 确保返回非空内容
            if response and response.strip():
                return response.strip()
            else:
                print("DEBUG: The response is empty, returning the original text")
                return text
            
        except Exception as e:
            print(f"LLM optimization failed: {str(e)}")
            # Ensure all possible exceptions are handled
            return text

    def _describe_image_with_llm(self, image_path: str) -> str:
        """When OCR fails, use LLM to describe the image content."""
        try:
            # Read image file and convert to base64
            with open(image_path, "rb") as image_file:
                image_data = base64.b64encode(image_file.read()).decode('utf-8')
            
            # Prepare messages
            messages = [
                {
                    "role": MessageRole.SYSTEM,
                    "content": [
                        {
                            "type": "text",
                            "text": "你是一个图像分析专家。请详细描述图像中的内容，特别关注任何文本、符号、图表或重要视觉元素。"
                        }
                    ],
                },
                {
                    "role": MessageRole.USER,
                    "content": [
                        {
                            "type": "text",
                            "text": "请详细描述这张图像中的内容，特别是任何可见的文本:"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{image_data}"
                            }
                        }
                    ],
                }
            ]
            
            # Call model
            response = self.model(messages).content
            
            if response and response.strip():
                return f"[图像描述] {response.strip()}"
            else:
                return "无法分析图像内容"
            
        except Exception as e:
            print(f"Image description failed: {str(e)}")
            return f"无法处理图像: {str(e)}"

class FileProcessingResult:
    """Standardized format for file processing results"""
    def __init__(self, title: Optional[str] = None, content: str = "", error: Optional[str] = None):
        self.title = title
        self.content = content
        self.error = error

    def __str__(self):
        if self.error:
            return f"Error: {self.error}"
        if self.title:
            return f"{self.title}\n\n{self.content}"
        return self.content

class FileProcessingException(Exception):
    """Base class for file processing exceptions"""
    pass

class UnsupportedFormatException(FileProcessingException):
    """Unsupported file format exception"""
    pass

class FileConversionException(FileProcessingException):
    """File conversion failed exception"""
    pass

class PDFTool(Tool):
    name = "pdf_tool"
    description = "Extract text from PDF files and perform analysis."
    inputs = {
        "file_path": {
            "type": "string",
            "description": "Path to the PDF file",
            "nullable": True
        },
        "page_range": {
            "type": "string", 
            "description": "Page range to extract (e.g., '1-5', 'all')",
            "default": "all",
            "nullable": True
        }
    }
    output_type = "string"

    def __init__(self, file_processor: FileProcessor):
        super().__init__()
        self.file_processor = file_processor
        
    def forward(self, file_path: str = None, page_range: str = "all") -> str:
        """Extract text from a PDF file."""
        if not file_path or not os.path.exists(file_path):
            return f"Error: File does not exist or path is invalid - {file_path}"
            
        try:
            print(f"Attempting to extract text from PDF: {file_path}")
            text = ""
            
            pages = None
            if page_range != "all":
                try:
                    if "-" in page_range:
                        start, end = map(int, page_range.split("-"))
                        pages = list(range(start-1, end))  # pdfminer使用0索引
                    else:
                        pages = [int(page_range) - 1]
                except ValueError:
                    return f"Error: Invalid page range - {page_range}"
            
            try:
                print("Using pdfminer to extract text...")
                text = pdfminer.high_level.extract_text(file_path, page_numbers=pages)
            except Exception as e:
                print(f"pdfminer extraction failed: {e}")
                return f"Error: PDF processing failed - {str(e)}"
                
            if text and len(text.strip()) > 50:
                print(f"Successfully extracted text, length: {len(text)} characters")
                return format_pdf_text(text, file_path)
                
            print("pdfminer extraction failed or text too short, trying PyPDF2...")
            import PyPDF2
            text = ""
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                if pages is None:
                    pages = range(len(reader.pages))
                for i in pages:
                    if i < len(reader.pages):
                        page_text = reader.pages[i].extract_text()
                        if page_text:
                            text += f"\n--- 第 {i+1} 页 ---\n" + page_text
            
            if text and len(text.strip()) > 50:
                print(f"PyPDF2 successfully extracted text, length: {len(text)} characters")
                return format_pdf_text(text, file_path)
                
            print("Direct text extraction failed, PDF may be a scanned document, trying OCR...")
            
            # Check if OCR tool is available
            if hasattr(self.file_processor, 'ocr_reader'):
                # Use PDF to image then OCR method
                from pdf2image import convert_from_path
                import tempfile
                
                print(f"Converting PDF to images for OCR processing...")
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Determine which pages to process
                    if pages is None:
                        pdf_pages = convert_from_path(file_path)
                    else:
                        pdf_pages = convert_from_path(file_path, first_page=min(pages)+1, last_page=max(pages)+1)
                    
                    all_text = []
                    for i, page in enumerate(pdf_pages):
                        # Save image
                        image_path = os.path.join(temp_dir, f'page_{i+1}.png')
                        page.save(image_path, 'PNG')
                        
                        # OCR processing
                        print(f"OCR processing page {i+1}...")
                        result = self.file_processor.ocr_reader.readtext(image_path)
                        page_text = "\n".join([text for _, text, _ in result])
                        all_text.append(f"\n--- Page {i+1} (OCR) ---\n{page_text}")
                    
                    if all_text:
                        ocr_text = "\n".join(all_text)
                        print(f"OCR successfully extracted text, length: {len(ocr_text)} characters")
                        return format_pdf_text(ocr_text, file_path, is_ocr=True)
            
            # If all methods fail
            return f"Warning: Unable to extract text from PDF. This PDF may be a scanned document, image-based document, or protected with special security features. Please upload a screenshot of the PDF for analysis."
            
        except Exception as e:
            import traceback
            trace = traceback.format_exc()
            return f"Error: PDF processing failed - {str(e)}\n\nDetailed error:\n{trace}"

def format_pdf_text(text, file_path, is_ocr=False):
    """Format extracted PDF text"""
    file_name = os.path.basename(file_path)
    method = "OCR" if is_ocr else "Direct extraction"
    
    # Get text statistics
    text_length = len(text)
    lines = text.count('\n') + 1
    words = len(text.split())
    
    header = f"""# PDF file analysis: {file_name}
- Extraction method: {method}
- Text length: {text_length} characters
- Number of lines: {lines}
- Number of words: {words}

## Extracted content:
"""
    
    return header + text

class DOCXTool(Tool):
    name = "DOCX_Tool"
    description = "Convert DOCX files to HTML."
    inputs = {
        "file_path": {
            "type": "string",
            "description": "Path to the DOCX file.",
            "nullable": True
        }
    }
    output_type = "string"

    def __init__(self, file_processor: FileProcessor):
        super().__init__()
        self.file_processor = file_processor

    def forward(self, file_path: str = None) -> str:
        try:
            if file_path is None:
                return "Error: No DOCX file path provided"
                
            if not os.path.exists(file_path):
                return f"Error: File not found at path: {file_path}"
                
            # Check file extension
            if not file_path.lower().endswith('.docx'):
                return "Error: File is not a DOCX document."
                
            with open(file_path, "rb") as docx_file:
                result = mammoth.convert_to_html(docx_file)
                html_content = result.value
                
                if not html_content.strip():
                    return "No content found in the document."
                    
                return html_content.strip()
        except Exception as e:
            return f"Error processing DOCX: {str(e)}"

class XLSXTool(Tool):
    name = "XLSX_Tool"
    description = "Convert XLSX files to Markdown."
    inputs = {
        "file_path": {
            "type": "string",
            "description": "Path to the XLSX file.",
            "nullable": True
        }
    }
    output_type = "string"

    def __init__(self, file_processor: FileProcessor):
        super().__init__()
        self.file_processor = file_processor

    def forward(self, file_path: str = None) -> str:
        try:
            if file_path is None:
                return "Error: No Excel file path provided"
                
            if not os.path.exists(file_path):
                return f"Error: File not found at path: {file_path}"
                
            # Check file extension
            if not file_path.lower().endswith(('.xlsx', '.xls')):
                return "Error: File is not an Excel spreadsheet."
                
            sheets = pd.read_excel(file_path, sheet_name=None)
            if not sheets:
                return "No data found in the Excel file."
                
            md_content = ""
            for sheet_name, sheet_data in sheets.items():
                if sheet_data.empty:
                    continue
                    
                md_content += f"## {sheet_name}\n"
                md_content += sheet_data.to_markdown(index=False) + "\n\n"
                
            if not md_content:
                return "No data found in the Excel file."
                
            return md_content.strip()
        except Exception as e:
            return f"Error processing XLSX: {str(e)}"

class PPTXTool(Tool):
    name = "pptx_tool"
    description = "Extract text and structure from PowerPoint presentations, including image analysis."
    inputs = {
        "file_path": {
            "type": "string",
            "description": "Path to the PPTX file",
            "nullable": True
        }
    }
    output_type = "string"

    def __init__(self, file_processor: FileProcessor):
        super().__init__()
        self.file_processor = file_processor
        self.image_analyzer = None
        
    def forward(self, file_path: str = None) -> str:
        try:
            if not file_path or not os.path.exists(file_path):
                return f"Error: File does not exist or path is invalid - {file_path}"
                
            # Verify file type
            if self.file_processor.detect_file_type(file_path) != '.pptx':
                return "Error: File is not a PPTX format"
                
            # Extract PPT content
            presentation = pptx.Presentation(file_path)
            content = []
            
            # Add title
            content.append(f"# PowerPoint: {os.path.basename(file_path)}")
            
            # Create temporary directory for saving images
            temp_dir = tempfile.mkdtemp()
            
            try:
                # Initialize image analyzer (if possible)
                self._init_image_analyzer()
                
                # Process each slide
                slide_num = 0
                for slide in presentation.slides:
                    slide_num += 1
                    content.append(f"\n## Slide {slide_num}")
                    
                    # Process slide title
                    title = slide.shapes.title
                    if title and title.has_text_frame:
                        content.append(f"### {title.text.strip()}")
                    
                    # Process various shapes in the slide
                    for shape in slide.shapes:
                        # Skip already processed title
                        if shape == title:
                            continue
                            
                        # Process picture
                        if self._is_picture(shape):
                            # Try to get alt text
                            alt_text = ""
                            try:
                                alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                            except Exception:
                                pass
                            
                            # Extract and save picture
                            image_path = self._extract_image(shape, temp_dir, slide_num)
                            if image_path:
                                # Use ImageAnalysisTool to analyze picture
                                image_description = self._analyze_image(image_path)
                                
                                # Add picture description
                                if alt_text:
                                    content.append(f"[Picture description: {alt_text}]")
                                content.append(f"[Picture analysis: {image_description}]")
                            else:
                                content.append("[Picture: Unable to extract]")
                        
                        # Process table
                        elif self._is_table(shape):
                            content.append("\n#### Table content:")
                            table_content = []
                            
                            # Header
                            if shape.table.rows:
                                header = []
                                for cell in shape.table.rows[0].cells:
                                    header.append(cell.text.strip())
                                table_content.append("| " + " | ".join(header) + " |")
                                
                                # Separator line
                                table_content.append("| " + " | ".join(["---"] * len(header)) + " |")
                                
                                # Table content
                                for row in shape.table.rows[1:]:
                                    row_content = []
                                    for cell in row.cells:
                                        row_content.append(cell.text.strip())
                                    table_content.append("| " + " | ".join(row_content) + " |")
                                    
                                content.append("\n".join(table_content))
                        
                        # Process text box
                        elif shape.has_text_frame:
                            text = shape.text.strip()
                            if text:
                                content.append(text)
                    
                    # Process notes
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            content.append("\n#### Notes:")
                            content.append(notes)
            
            finally:
                # Clean up temporary directory
                shutil.rmtree(temp_dir, ignore_errors=True)
                
            return "\n\n".join(content)
            
        except Exception as e:
            return f"Error: PPTX processing failed - {str(e)}"
    
    def _init_image_analyzer(self):
        """Initialize image analyzer"""
        if self.image_analyzer is not None:
            return  # Already initialized
            
        try:
            # Check if FileProcessor has a model
            if hasattr(self.file_processor, 'model') and self.file_processor.model is not None:
                # Create image analyzer
                self.image_analyzer = ImageAnalysisTool(self.file_processor, self.file_processor.model)
            else:
                print("Warning: FileProcessor does not have a usable model, cannot perform image analysis")
        except Exception as e:
            print(f"Failed to initialize image analyzer: {str(e)}")
    
    def _is_picture(self, shape):
        """Check if shape is a picture"""
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            try:
                if hasattr(shape, "image"):
                    return True
            except:
                pass
        return False
        
    def _is_table(self, shape):
        """Check if shape is a table"""
        return hasattr(shape, "table")
        
    def _extract_image(self, shape, temp_dir, slide_num):
        """Extract image and save to temporary directory"""
        try:
            # Extract image
            if hasattr(shape, "image"):
                # Get image from shape
                image_bytes = shape.image.blob
                image_ext = shape.image.ext
                image_filename = f"slide_{slide_num}_{shape.name}.{image_ext}"
                image_path = os.path.join(temp_dir, image_filename)
                
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                    
                return image_path
            return None
        except Exception as e:
            print(f"Image extraction failed: {str(e)}")
            return None
            
    def _analyze_image(self, image_path):
        """Use ImageAnalysisTool to analyze image"""
        try:
            # Check if image analyzer is available
            if self.image_analyzer is None:
                return "Unable to analyze image: Image analyzer not initialized"
                
            # Use ImageAnalysisTool to analyze image
            result = self.image_analyzer.forward(image_path)
            
            # Simplify result (remove title, etc.)
            if result and isinstance(result, str):
                # Remove possible title and prefix
                result = re.sub(r'^.*?\[图像描述\]\s*', '', result, flags=re.DOTALL)
                return result.strip()
            
            return "Failed to analyze image content"
            
        except Exception as e:
            print(f"Image analysis failed: {str(e)}")
            return f"Image analysis failed: {str(e)}"

class ImageAnalysisTool(Tool):
    name = "Image_Analysis_Tool"
    description = "Analyze image content, provide detailed image descriptions. Input image path, output comprehensive analysis of the image."
    inputs = {
        "image_path": {
            "type": "string",
            "description": "Path to the image file.",
            "nullable": True
        }
    }
    output_type = "string"

    def __init__(self, file_processor: FileProcessor, model):
        """
        Initialize image analysis tool
        
        Args:
            file_processor: File processor instance
            model: Model for analyzing images
        """
        super().__init__()
        self.file_processor = file_processor
        self.model = model

    def forward(self, image_path: str = None) -> str:
        """
        Analyze image and return model-generated description
        
        Args:
            image_path: Image file path
        
        Returns:
            str: Detailed description of image content
        """
        try:
            if image_path is None:
                return "Error: No image file path provided"
                
            # Check if file exists
            if not os.path.exists(image_path):
                return f"Error: File {image_path} does not exist"
                
            # Get file extension
            ext = os.path.splitext(image_path)[1].lower()
            
            # Check if it is a supported image format
            if ext not in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.webp']:
                return f"Error: Unsupported image format {ext}"
                
            # Read image file and encode to base64
            with open(image_path, "rb") as image_file:
                encoded_image = base64.b64encode(image_file.read()).decode('utf-8')
                
            # Build message content
            messages = [
                {
                    "role": "system",
                    "content": [
                        {
                            "type": "text",
                            "text": "你是一位专业的图像分析助手。请详细描述图像中的内容，包括可见的物体、场景、人物特征、文字、行为、背景和任何其他重要细节。"
                        }
                    ],
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "请详细分析这张图片的内容:"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/{ext[1:]};base64,{encoded_image}"
                            }
                        }
                    ],
                }
            ]
            
            # 调用模型进行图像分析
            response = self.model(messages)
            
            # 返回分析结果
            return response.content if hasattr(response, 'content') else str(response)
                
        except Exception as e:
            return f"Error analyzing image: {str(e)}"





